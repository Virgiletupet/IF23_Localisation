"""
27_make_pptx.py — Support de soutenance IF23 P, thème UTT (16:9).
Logo + bandeau UTT, icônes Lucide en pastilles, sommaire paginé, pieds de page
numérotés, figures réutilisées. Sortie : reports/Soutenance_IF23_ALT3.pptx
"""
from __future__ import annotations

import sys
from pathlib import Path

try:
    sys.stdout.reconfigure(encoding="utf-8")
except Exception:
    pass

from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

ROOT = Path(__file__).resolve().parents[1]
FIG = ROOT / "reports" / "figures"
AST = ROOT / "reports" / "assets"
ICO = AST / "icons"
OUT = ROOT / "reports" / "Soutenance_IF23_ALT3.pptx"

NAVY = RGBColor(0x00, 0x1E, 0x62); CYAN = RGBColor(0x00, 0x9C, 0xDE)
BLUE = RGBColor(0x1F, 0x77, 0xB4); LIGHT = RGBColor(0xE3, 0xEE, 0xF9)
GREY = RGBColor(0x59, 0x59, 0x59); RED = RGBColor(0xC0, 0x39, 0x2B)
GREEN = RGBColor(0x2E, 0x8B, 0x57); WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT = "Segoe UI"
EMU = 914400

prs = Presentation()
prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def slide():
    return prs.slides.add_slide(BLANK)


def _set(run, size, color, bold=False, italic=False):
    run.font.name = FONT; run.font.size = Pt(size); run.font.color.rgb = color
    run.font.bold = bold; run.font.italic = italic


def textbox(s, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = Pt(2); tf.margin_right = Pt(2); tf.margin_top = Pt(1); tf.margin_bottom = Pt(1)
    return tf


def line(tf, text, size, color, bold=False, italic=False, align=PP_ALIGN.LEFT, first=False, space=6):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align; p.space_after = Pt(space)
    r = p.add_run(); r.text = text; _set(r, size, color, bold, italic)
    return p


def rrect(s, x, y, w, h, fill, linecol=None):
    sp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.shadow.inherit = False; sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if linecol is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = linecol; sp.line.width = Pt(1)
    return sp


def rect(s, x, y, w, h, fill):
    sp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.shadow.inherit = False; sp.fill.solid(); sp.fill.fore_color.rgb = fill; sp.line.fill.background()
    return sp


def oval(s, x, y, d, fill):
    sp = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    sp.shadow.inherit = False; sp.fill.solid(); sp.fill.fore_color.rgb = fill; sp.line.fill.background()
    return sp


def img_fit(s, path, x, y, max_w, max_h, center_x=None):
    p = Path(path)
    if not p.exists():
        return None
    iw, ih = Image.open(p).size
    ratio = min((max_w * EMU) / iw, (max_h * EMU) / ih)
    w = int(iw * ratio); h = int(ih * ratio)
    px = Inches(x) if center_x is None else Inches(center_x) - w // 2
    return s.shapes.add_picture(str(p), px, Inches(y), width=Emu(w), height=Emu(h))


def icon_pastille(s, icon, x, y, d=0.62, fill=LIGHT):
    oval(s, x, y, d, fill)
    ip = ICO / f"{icon}.png"
    if ip.exists():
        pad = d * 0.26
        img_fit(s, ip, x + pad, y + pad, d - 2 * pad, d - 2 * pad)


def title_bar(s, title, icon=None):
    rect(s, 0, 0, SW.inches, 1.05, NAVY)
    rect(s, 0, 1.05, SW.inches, 0.06, CYAN)
    tx = 0.55
    if icon:
        oval(s, 0.45, 0.22, 0.62, RGBColor(0x12, 0x34, 0x78))
        ip = ICO / f"{icon}.png"
        if ip.exists():
            img_fit(s, ip, 0.45 + 0.16, 0.22 + 0.16, 0.30, 0.30)
        tx = 1.30
    tf = textbox(s, tx, 0.0, 11.4, 1.05, anchor=MSO_ANCHOR.MIDDLE)
    line(tf, title, 26, WHITE, bold=True, first=True, space=0)


def bullets(s, items, x, y, w, h, size=18, gap=10):
    tf = textbox(s, x, y, w, h)
    for i, it in enumerate(items):
        if isinstance(it, tuple):
            pre, rest = it
        else:
            pre, rest = None, it
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        b = p.add_run(); b.text = "▪  "; _set(b, size, CYAN, bold=True)
        if pre:
            r0 = p.add_run(); r0.text = pre; _set(r0, size, NAVY, bold=True)
        r = p.add_run(); r.text = rest; _set(r, size, GREY)
    return tf


def kpi(s, x, y, w, h, value, label, color=NAVY):
    rrect(s, x, y, w, h, LIGHT)
    tf = textbox(s, x, y + 0.10, w, h - 0.2, anchor=MSO_ANCHOR.MIDDLE)
    line(tf, value, 30, color, bold=True, align=PP_ALIGN.CENTER, first=True, space=2)
    line(tf, label, 12.5, GREY, align=PP_ALIGN.CENTER, space=0)


# ---------------------------------------------------------------- diapos
def s_title():
    s = slide()
    img_fit(s, AST / "utt_banner.png", 0, 0, 13.34, 2.5, center_x=SW.inches / 2)
    rect(s, 0, 2.46, SW.inches, 0.06, CYAN)
    img_fit(s, AST / "utt_logo.png", 0, 2.95, 2.6, 1.0, center_x=SW.inches / 2)
    tf = textbox(s, 1.0, 3.95, 11.33, 2.6, anchor=MSO_ANCHOR.TOP)
    line(tf, "Localisation indoor par fingerprinting WiFi", 34, NAVY, bold=True,
         align=PP_ALIGN.CENTER, first=True, space=6)
    line(tf, "Détection de salle et estimation de position dans le bâtiment de l'UTT", 18, CYAN,
         align=PP_ALIGN.CENTER, space=18)
    line(tf, "Binôme 6  —  Molka GHADDAB · Maxime LE GAL · Virgile TUPET", 16, NAVY, bold=True,
         align=PP_ALIGN.CENTER, space=4)
    line(tf, "Encadrante : Farah Chehade   |   UV IF23 P   |   Soutenance — 3 juillet 2026", 13,
         GREY, align=PP_ALIGN.CENTER, space=0)


def s_sommaire():
    s = slide(); title_bar(s, "Sommaire", "carte")
    items = [("Contexte et objectifs", 3), ("Données et caractéristiques", 4),
             ("Classification et évaluation réaliste", 6), ("Évolution temporelle (drift)", 8),
             ("Régression intra-zone (x, y)", 9), ("Cartographie", 10),
             ("Difficultés et conclusion", 12)]
    icons = ["wifi", "donnees", "chute", "horloge", "regle", "boussole", "fusee"]
    y = 1.6
    for (txt, n), ic in zip(items, icons):
        icon_pastille(s, ic, 1.4, y, 0.6)
        tf = textbox(s, 2.2, y, 8.8, 0.6, anchor=MSO_ANCHOR.MIDDLE)
        line(tf, txt, 19, NAVY, bold=True, first=True, space=0)
        tf2 = textbox(s, 11.0, y, 1.4, 0.6, anchor=MSO_ANCHOR.MIDDLE)
        line(tf2, f"diapo {n}", 14, CYAN, align=PP_ALIGN.RIGHT, first=True, space=0)
        y += 0.74


def s_contexte():
    s = slide(); title_bar(s, "Contexte et objectifs", "wifi")
    bullets(s, [
        ("Problème : ", "à l'intérieur, le GPS ne fonctionne pas (murs, étages)."),
        ("Idée : ", "le WiFi est partout ; chaque lieu a une « empreinte » de puissances (RSSI)."),
        ("Fingerprinting : ", "on enregistre les empreintes (hors ligne), puis on compare une mesure (en ligne)."),
    ], 0.6, 1.5, 7.4, 3.2, size=18, gap=14)
    # deux objectifs en cartes
    rrect(s, 0.6, 4.7, 5.9, 2.2, LIGHT)
    icon_pastille(s, "cible", 0.85, 4.95, 0.6, WHITE)
    tf = textbox(s, 1.6, 4.95, 4.7, 1.9)
    line(tf, "Classification", 18, NAVY, bold=True, first=True, space=2)
    line(tf, "Prédire la zone (salle / couloir) parmi N", 14, GREY, space=0)
    rrect(s, 6.7, 4.7, 5.9, 2.2, LIGHT)
    icon_pastille(s, "regle", 6.95, 4.95, 0.6, WHITE)
    tf = textbox(s, 7.7, 4.95, 4.7, 1.9)
    line(tf, "Régression", 18, NAVY, bold=True, first=True, space=2)
    line(tf, "Estimer la position (x, y) dans une zone", 14, GREY, space=0)
    img_fit(s, FIG / "fig22_plan_etage.png", 8.2, 1.45, 4.9, 3.0)


def s_donnees():
    s = slide(); title_bar(s, "Données : trois campagnes", "donnees")
    cards = [("ALT1", "février 2026", "24 zones · 1 498 scans", "apprentissage", BLUE),
             ("ALT2", "28 avril 2026", "10 zones · 452 scans", "test temporel", CYAN),
             ("ALT3", "juin 2026", "à collecter", "test final", GREY)]
    x = 0.6
    for code, per, vol, role, col in cards:
        rrect(s, x, 1.5, 3.95, 2.3, LIGHT)
        rect(s, x, 1.5, 3.95, 0.12, col)
        tf = textbox(s, x + 0.25, 1.75, 3.5, 2.0)
        line(tf, code, 24, NAVY, bold=True, first=True, space=2)
        line(tf, per, 14, col, bold=True, space=6)
        line(tf, vol, 14, GREY, space=2)
        line(tf, role, 13, GREY, italic=True, space=0)
        x += 4.15
    bullets(s, [
        ("Mesure : ", "SSID, BSSID, RSSI (dBm), horodatage ; vecteur de puissances, -100 si absent."),
        ("Régression : ", "une salle (5 × 7 m) cartographiée finement — 46 points (x, y) × 52 bornes."),
        ("Objet : ", "mesurer la tenue du modèle dans le temps (stratégies ALT1→ALT2→ALT3)."),
    ], 0.6, 4.2, 12.1, 2.6, size=17, gap=12)


def s_features():
    s = slide(); title_bar(s, "Caractéristiques et choix de l'identifiant", "couches")
    bullets(s, [
        ("Vectorisation : ", "statistiques de RSSI par réseau (moyenne, max, présence) + indicateurs globaux."),
        ("Sélection : ", "285 paires de réseaux très corrélées → on évite la redondance."),
    ], 0.6, 1.5, 12.1, 1.6, size=18, gap=12)
    # 3 cartes encodage
    enc = [("SSID seul", "13 réseaux", "98,3 %", BLUE), ("BSSID seul", "104 réseaux", "99,0 %", CYAN),
           ("SSID | BSSID", "139 réseaux", "100 %", GREEN)]
    x = 0.9
    for name, n, acc, col in enc:
        rrect(s, x, 3.4, 3.7, 2.6, LIGHT)
        tf = textbox(s, x + 0.2, 3.65, 3.3, 2.2, anchor=MSO_ANCHOR.MIDDLE)
        line(tf, name, 19, NAVY, bold=True, align=PP_ALIGN.CENTER, first=True, space=4)
        line(tf, n, 14, GREY, align=PP_ALIGN.CENTER, space=10)
        line(tf, acc, 30, col, bold=True, align=PP_ALIGN.CENTER, space=2)
        line(tf, "accuracy hold-out", 11, GREY, align=PP_ALIGN.CENTER, space=0)
        x += 4.05
    tf = textbox(s, 0.6, 6.25, 12.1, 0.7)
    line(tf, "L'identifiant composé SSID|BSSID apporte le maximum d'information (résultat de hold-out).",
         14, GREY, italic=True, first=True, space=0)


def s_classif():
    s = slide(); title_bar(s, "Classification : l'importance d'une évaluation réaliste", "chute")
    img_fit(s, FIG / "fig13_escalier.png", 0.5, 1.4, 8.2, 5.2)
    kpi(s, 9.0, 1.7, 3.7, 1.4, "99 % → 47 %", "test aléatoire → test réel ALT2", RED)
    bullets(s, [
        "Un découpage aléatoire surestime (mesures quasi identiques).",
        "On durcit le protocole : séparation temporelle, puis nouvelle campagne.",
        "La performance réaliste rejoint ce qui est observé sur le terrain.",
    ], 9.0, 3.3, 3.9, 3.4, size=15, gap=12)


def s_importance():
    s = slide(); title_bar(s, "Sur quoi le modèle s'appuie", "cpu")
    img_fit(s, FIG / "fig4_importance_reseaux.png", 0.5, 1.4, 7.6, 5.3)
    kpi(s, 8.5, 1.7, 4.2, 1.4, "99,6 %", "bonne zone dans le Top-3", GREEN)
    bullets(s, [
        ("Constat : ", "les 2 réseaux les plus discriminants sont des téléphones personnels."),
        ("Conséquence : ", "utiles sur le moment, absents la campagne suivante."),
        ("Top-3 fiable : ", "idéal pour proposer une courte liste de zones."),
    ], 8.5, 3.3, 4.4, 3.4, size=15, gap=12)


def s_evolution():
    s = slide(); title_bar(s, "Évolution temporelle : le « drift » du WiFi", "horloge")
    img_fit(s, FIG / "fig18_evolution.png", 0.5, 1.45, 8.3, 4.6)
    kpi(s, 9.1, 1.7, 1.85, 1.5, "47 %", "ALT1 → ALT2", RED)
    kpi(s, 11.05, 1.7, 1.85, 1.5, "100 %", "+ ré-app.", GREEN)
    bullets(s, [
        "Modèle de février testé sur avril : la moitié des zones décroche.",
        "Cause : l'environnement radio change (bornes, réseaux invités).",
        "Réintégrer 50 % d'ALT2 restaure la performance.",
        ("À retenir : ", "ré-entraîner régulièrement."),
    ], 9.1, 3.5, 3.9, 3.2, size=14.5, gap=10)


def s_regression():
    s = slide(); title_bar(s, "Régression intra-zone (position x, y)", "regle")
    img_fit(s, FIG / "fig20_regression_scatter.png", 0.6, 1.5, 4.5, 5.0)
    img_fit(s, FIG / "fig19_regression_cdf.png", 5.1, 1.7, 4.2, 3.2)
    kpi(s, 9.6, 1.7, 3.3, 1.4, "≈ 1,1 m", "erreur médiane (zone 5×7 m)", NAVY)
    bullets(s, [
        ("Modèle hybride : ", "X par processus gaussien, Y par ExtraTrees."),
        ("Pourquoi : ", "les deux axes n'ont pas la même difficulté."),
        ("Limite : ", "précision bornée par le nombre de points (46)."),
    ], 9.6, 3.3, 3.5, 3.4, size=14.5, gap=11)


def s_carto1():
    s = slide(); title_bar(s, "Cartographie : puissance et plan", "boussole")
    img_fit(s, FIG / "fig21_heatmaps.png", 0.5, 1.45, 6.2, 5.4)
    img_fit(s, FIG / "fig22_plan_etage.png", 6.9, 1.7, 6.1, 3.7)
    bullets(s, [
        ("Carte de puissance : ", "RSSI interpolé par borne sur la zone."),
        ("Plan : ", "campus UTT (bâtiment) + croquis d'étage (salle, position)."),
    ], 6.9, 5.4, 6.1, 1.6, size=15, gap=10)


def s_carto2():
    s = slide(); title_bar(s, "Cartographie : démo et aide à la mobilité", "route")
    img_fit(s, FIG / "fig24_demo_localisation.png", 0.5, 1.5, 7.7, 4.0)
    img_fit(s, FIG / "fig23_map_assist.png", 8.4, 1.6, 4.5, 3.2)
    kpi(s, 8.7, 5.0, 3.9, 1.2, "+7,6 pts", "aide mobilité (désactivable)", GREEN)
    bullets(s, [
        "Scan réel → étage, zone (Top-3), position sur le plan.",
        "Aide mobilité : graphe de transition entre zones voisines.",
        ("Prudence : ", "activable, désactivée par défaut."),
    ], 0.6, 5.7, 7.7, 1.4, size=14.5, gap=8)


def s_difficultes():
    s = slide(); title_bar(s, "Difficultés rencontrées", "alerte")
    bullets(s, [
        ("Encodage live : ", "format SSID|BSSID attendu par le modèle ≠ SSID brut du scanner → détection auto du format."),
        ("Renouvellement de l'infra : ", "des bornes remplacées rendent invisibles les empreintes connues (limite du fingerprinting)."),
        ("Petites classes & versions : ", "zone trop peu mesurée écartée de la validation ; versions scikit-learn gérées."),
        ("Protocole de mesure : ", "passage du découpage A/B (ALT1) à la capture continue (ALT2)."),
    ], 0.7, 1.7, 12.0, 4.8, size=18, gap=20)


def s_conclusion():
    s = slide(); title_bar(s, "Conclusion et perspectives", "fusee")
    bullets(s, [
        ("Deux tâches livrées : ", "classification de zone et position (x, y) à ≈ 1 m."),
        ("Enseignement clé : ", "le hold-out flatte ; le vrai test (nouvelle campagne) tombe à 47 % puis se restaure."),
        ("Cartographie : ", "puissance, plan, démo bout-en-bout, aide mobilité désactivable."),
    ], 0.7, 1.6, 12.0, 2.7, size=18, gap=14)
    line(textbox(s, 0.7, 4.35, 12, 0.5), "Perspectives", 19, NAVY, bold=True, first=True, space=0)
    bullets(s, [
        "Collecter ALT3 (juin) → stratégie 3.",
        "Densifier la régression et l'étendre à d'autres zones.",
        "Intégrer la cartographie dans l'interface temps réel.",
    ], 0.7, 4.85, 12.0, 2.0, size=16, gap=8)


def s_merci():
    s = slide()
    rect(s, 0, 0, SW.inches, SH.inches, NAVY)
    rect(s, 0, 3.55, SW.inches, 0.06, CYAN)
    img_fit(s, AST / "utt_logo.png", 0, 1.5, 2.6, 1.0, center_x=SW.inches / 2)
    tf = textbox(s, 1, 3.7, 11.33, 2.5, anchor=MSO_ANCHOR.TOP)
    line(tf, "Merci de votre attention", 34, WHITE, bold=True, align=PP_ALIGN.CENTER, first=True, space=8)
    line(tf, "Questions ?", 22, CYAN, align=PP_ALIGN.CENTER, space=14)
    line(tf, "Binôme 6 — Molka GHADDAB · Maxime LE GAL · Virgile TUPET", 14, WHITE,
         align=PP_ALIGN.CENTER, space=0)


def footer_numbers():
    total = len(prs.slides._sldIdLst)
    for i, sl in enumerate(prs.slides, 1):
        if i == 1 or i == total:   # ni titre ni clôture
            continue
        img_fit(sl, AST / "utt_logo.png", 0.35, 7.0, 1.0, 0.38)
        tf = textbox(sl, 1.5, 7.0, 9.0, 0.4, anchor=MSO_ANCHOR.MIDDLE)
        line(tf, "IF23 P — Localisation indoor WiFi", 9, GREY, first=True, space=0)
        tf2 = textbox(sl, 11.4, 7.0, 1.6, 0.4, anchor=MSO_ANCHOR.MIDDLE)
        line(tf2, f"{i} / {total}", 10, GREY, align=PP_ALIGN.RIGHT, first=True, space=0)


def main():
    s_title(); s_sommaire(); s_contexte(); s_donnees(); s_features()
    s_classif(); s_importance(); s_evolution(); s_regression()
    s_carto1(); s_carto2(); s_difficultes(); s_conclusion(); s_merci()
    footer_numbers()
    OUT.parent.mkdir(parents=True, exist_ok=True); prs.save(OUT)
    print("PPTX généré :", OUT, "|", len(prs.slides._sldIdLst), "diapos")


if __name__ == "__main__":
    main()
